#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0, 82, 147)
DARK_BLUE = RGBColor(0, 51, 102)
GRAY = RGBColor(128, 128, 128)
LIGHT_BLUE = RGBColor(0, 102, 204)
ORANGE = RGBColor(255, 128, 0)

def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Cover
add_title_slide(prs, 
    "AI Agent 研发实战深度总结",
    "QCon 2026 北京 · Agentic Engineering & Coding Agent")

# Slide 2: 核心转变
add_content_slide(prs, 
    "为什么这是范式级别的变革",
    [
        "代码第一次从「思考载体」变成「执行载体」（黄东旭）",
        "模型 = CPU，Agent = OS，Skills = App",
        "入口迁移：GUI → IM → 自然语言",
        "协作模式：从「人找事」到「事找人」",
        "2026分水岭：从试点热潮到规模落地"
    ])

# Slide 3: 企业级Agent架构
add_two_column_slide(prs,
    "企业级 Agent 架构核心要素",
    "六个核心技术要素",
    [
        "Context Window：单次推理的工作记忆上限",
        "System Prompt：定义Agent身份和角色",
        "Skills：领域知识包，按需加载",
        "MCP：连接外部系统的标准协议",
        "Memory：跨会话的持久化信息",
        "Sub-agent：委派子任务"
    ],
    "ZillizTown 企业实践",
    [
        "飞书IM入口：Avatar住在IM里",
        "Cron自醒：Agent自己设闹钟巡检",
        "A2A协作：Avatar找Avatar帮忙",
        "对话训练：手机就能训练Avatar",
        "Hub路由：按Skills找对的Agent",
        "多Session：绕过1M上限"
    ])

# Slide 4: 安全五层模型
add_content_slide(prs,
    "Agent 安全五层攻击面模型（Sunny Duan）",
    [
        "L5 输入层：Prompt注入、多模态注入、Jailbreak绕过",
        "L4 规划层：目标劫持、规划操纵、幻觉利用",
        "L3 记忆层：记忆投毒、上下文窗口耗尽、RAG投毒",
        "L2 工具层：供应链攻击、MCP参数注入",
        "L1 反馈层：伪造反馈、错误抑制、观测遮蔽",
        "防御思路：用AI防御AI，以智能对抗智能"
    ])

# Slide 5: 三位一体编码
add_two_column_slide(prs,
    "Coding Agent 三位一体编码方案（邓立山）",
    "Rules（最高宪法）",
    [
        "目录结构约束",
        "架构模式约束", 
        "技术栈约束",
        "编码风格约束",
        "确保架构不偏离"
    ],
    "Skills（执法机构）",
    [
        "自动路由到不同场景",
        "工作流编排",
        "质量检查",
        "按需加载用完释放",
        "让方案自动运转"
    ])

# Slide 6: PRD到上线闭环
add_content_slide(prs,
    "从 PRD 到上线的交付闭环（李文鹏）",
    [
        "MVP文档基座 + 自然生长：成本低，不改工作方式",
        "Spec不是文档，是需求契约",
        "Hybrid Execution Protocol：Agent总成本=准备成本+验证成本+返工风险",
        "缺口门禁：Agent执行前的自检，把「先动手」变成「先确认」",
        "AI提的PR审核时间是人的4.6倍，平均打回轮次3→1.5",
        "公共模块变更导致的线上bug数：0"
    ])

# Slide 7: 飞轮模式
add_two_column_slide(prs,
    "Coding Agent 飞轮模式（牛万鹏）",
    "Feedback Loop",
    [
        "工具调用：Query/Tool Call比例",
        "上下文：Skills Tokens消耗",
        "执行结果：单Query更��文件数",
        "执行轨迹：Debug类任务路径",
        "MCP动态加载节省98% Tokens",
        "GPT偏好Bash而非专用工具"
    ],
    "Benchmark",
    [
        "从Git挖掘评测集",
        "另一Agent验证结果",
        "四象限分析：Outcome/Execution",
        "Outcome 60%：功能正确性、完整性",
        "Execution 40%：理解能力、执行质量",
        "发现异常值才能真正优化"
    ])

# Slide 8: Token成本优化
add_content_slide(prs,
    "DeepResearch Token 成本优化（卢亿雷）",
    [
        "Token是智能时代的「电力」：没有Token，AI停摆",
        "单次深度研究Token消耗：80K-300K tokens",
        "研究阶段占75%成本：多轮检索上下文膨胀",
        "Output单价是Input的3-5倍，控制输出长度ROI最高",
        "请求前：语义缓存 + Prompt Cache",
        "请求中：模型路由 + 输入压缩 + Output精确控制",
        "流程控制：无显著新信息时提前终止"
    ])

# Slide 9: 关键数据
add_content_slide(prs,
    "关键数据与实证",
    [
        "AI编码率：9.6% → 89.2%，提升9倍+（邓立山）",
        "Bug率：<0.2%，发布回滚率：<2%",
        "人均Query次数增长5倍（百度Comate）",
        "Tokens缓存命中率：80%+",
        "MCP动态加载节省：98%",
        "PR审核效率提升：4.6倍",
        "产出效率提升：10倍（20人团队）"
    ])

# Slide 10: 核心价值
add_content_slide(prs,
    "核心价值总结",
    [
        "范式转变：AI不是银弹，是超级杠杆",
        "引擎越强，方向盘越重要",
        "代码即规范，规范即代码",
        "Rules + Spec + Skills = 可复用的AI编码方案",
        "Feedback Loop + Benchmark = Agent行为可观测",
        "人也是Tool：全员Agent Engineer转型"
    ])

# Slide 11: 学习路径
add_content_slide(prs,
    "内部参考与学习路径",
    [
        "入门概念：黄东旭《自主系统的时代》、韦韬《智能革命时代》",
        "工程方法：邓立山《AI Coding全栈实战》、李文鹏《PRD到上线闭环》",
        "平台实践：牛万鹏《飞轮模式》、徐翔《JoyCode》",
        "安全专题：Sunny Duan《Agent安全》",
        "企业级落地：陈彪《ZillizTown》",
        "成本优化：卢亿雷《DeepResearch Token优化》"
    ])

# Save
output_path = "/Users/lidan/Desktop/软件大会/我的模块/AI-Agent-研发实战总结.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")